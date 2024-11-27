import json
from pathlib import Path

from config import Config
from docx import Document
from pptx import Presentation
from unstructured.documents.elements import Element
from unstructured.partition.docx import partition_docx
from unstructured.partition.pdf import partition_pdf
from unstructured.partition.pptx import partition_pptx
from unstructured.partition.xlsx import partition_xlsx

config = Config()


def create_file_output_folder(output_folder: Path, base_name: str) -> tuple[Path, Path]:
    """
    Create a specific folder for each file's outputs
    """
    file_output_folder = Path(output_folder) / base_name
    images_folder = file_output_folder / "images"
    file_output_folder.mkdir(parents=True, exist_ok=True)
    images_folder.mkdir(parents=True, exist_ok=True)
    return file_output_folder, images_folder


def docx_extract_images(doc_path: Path, images_folder: Path) -> list:
    doc = Document(doc_path)
    images = []
    image_count = 0
    base_name = doc_path.stem

    # Ensure the images_folder exists
    images_folder.mkdir(parents=True, exist_ok=True)

    # Extract all images
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            img = rel.target_part.blob
            image_path = images_folder / f"{base_name}_image{image_count}.jpg"
            with image_path.open("wb") as img_file:
                img_file.write(img)
            image_count += 1
            images.append(str(image_path))

    return images


def pptx_extract_images(pptx_path: Path, images_folder: Path) -> list:
    prs = Presentation(pptx_path)
    images = []
    image_count = 0
    base_name = pptx_path.stem

    # Ensure the images_folder exists
    images_folder.mkdir(parents=True, exist_ok=True)

    # Extract images from shapes in each slide
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                image = shape.image
                image_bytes = image.blob
                image_path = images_folder / f"{base_name}_image{image_count}.{image.ext}"
                with image_path.open("wb") as img_file:
                    img_file.write(image_bytes)
                image_count += 1
                images.append(str(image_path))

    return images


def extract_document_elements(
    input_folder: Path, images_folder: Path, fname: str
) -> list[Element] | None:
    """
    Extract elements from various document types.
    input_folder: Folder containing the document files
    images_folder: Folder to save extracted images
    fname: File name of the document
    """
    file_path = input_folder / fname

    common_params = {
        "filename": str(file_path),
        "strategy": "hi_res",
        "image_output_dir_path": str(images_folder),
    }

    if fname.lower().endswith(".pdf"):
        return partition_pdf(
            **common_params,
            extract_images_in_pdf=True,
            infer_table_structure=True,
            chunking_strategy="by_title",
            max_characters=4000,
            new_after_n_chars=3800,
            combine_text_under_n_chars=2000,
        )

    elif fname.lower().endswith((".pptx", ".ppt")):
        return partition_pptx(
            **common_params,
            include_page_breaks=True,
            extract_images_in_pptx=True,
        )

    elif fname.lower().endswith((".docx", ".doc")):
        return partition_docx(
            **common_params,
            extract_images_in_doc=True,
            infer_table_structure=True,
        )

    elif fname.lower().endswith((".xlsx", ".xlsm")):
        return partition_xlsx(**common_params, extract_tables=True)

    return None


def categorize_elements(raw_elements: list[Element]) -> tuple[list[str], list[str]]:
    """
    Categorize extracted elements from documents into tables and texts.
    raw_elements: List of unstructured.documents.elements
    """
    tables = []
    texts = []
    for element in raw_elements:
        element_type = str(type(element))
        if "Table" in element_type:
            tables.append(str(element))
        elif any(
            text_type in element_type
            for text_type in ["CompositeElement", "Text", "Title", "NarrativeText"]
        ):
            texts.append(str(element))
    return texts, tables


def process_documents(input_folder: Path, output_folder: Path) -> None:
    """
    Process all supported document types in the input folder
    """
    supported_extensions = (".pdf", ".pptx", ".ppt", ".docx", ".doc", ".xlsx", ".xlsm")

    for fname in input_folder.iterdir():
        if fname.suffix.lower() in supported_extensions:
            try:
                print(f"Processing {fname}...")
                base_name = fname.stem

                # Create specific folders for this file's outputs
                file_output_folder, images_folder = create_file_output_folder(
                    output_folder, base_name
                )

                if fname.suffix.lower() in (".docx", ".doc"):
                    # Extract images separately
                    _ = docx_extract_images(fname, images_folder)
                elif fname.suffix.lower() in (".pptx", ".ppt"):
                    # Extract images separately
                    _ = pptx_extract_images(fname, images_folder)
                # Extract elements with the specific images folder
                raw_elements = extract_document_elements(input_folder, images_folder, fname.name)
                texts, tables = categorize_elements(raw_elements)

                # Save texts
                text_file = file_output_folder / "texts.json"
                with text_file.open("w", encoding="utf-8") as f:
                    json.dump(texts, f, indent=4, ensure_ascii=False)

                # Save tables if any exist
                if tables:
                    table_file = file_output_folder / "tables.json"
                    with table_file.open("w", encoding="utf-8") as f:
                        json.dump(tables, f, indent=4, ensure_ascii=False)

                print(f"Successfully processed {fname}")

            except Exception as e:
                print(f"Error processing {fname}: {str(e)}")


def extract_advanced_data_demo() -> None:
    input_folder = config.RAW_DATA_FOLDER
    output_folder = config.PROCESSED_DATA_FOLDER
    input_folder.mkdir(parents=True, exist_ok=True)
    output_folder.mkdir(parents=True, exist_ok=True)
    process_documents(input_folder, output_folder)
